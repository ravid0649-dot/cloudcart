import subprocess
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

TITLE_BLUE = RGBColor(39, 95, 170)
BODY_BLACK = RGBColor(0, 0, 0)
ACCENT_RED = RGBColor(196, 30, 58)
LIGHT_BG = RGBColor(255, 255, 255)
LOGO_PATH = Path(
    "/Users/ravidesai/.cursor/projects/Users-ravidesai-Documents-cc2/assets/"
    "Screenshot_2026-04-21_at_10.17.15_AM-37896885-0bba-42b5-9f5f-2bbf461456bd.png"
)


def apply_theme(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(0.88), Inches(12.6), Inches(0.02))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = ACCENT_RED
    line_shape.line.fill.background()


def add_logo(slide) -> None:
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(9.25), Inches(0.06), width=Inches(3.35))


def style_title(title_frame, text: str, size: int = 34) -> None:
    title_frame.text = text
    p = title_frame.paragraphs[0]
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = TITLE_BLUE


def add_title_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_theme(slide)
    add_logo(slide)

    title_box = slide.shapes.title
    title_box.left = Inches(0.55)
    title_box.top = Inches(1.0)
    title_box.width = Inches(8.3)
    title_box.height = Inches(0.8)
    style_title(title_box.text_frame, title, size=35)

    subtitle_placeholder = slide.placeholders[1]
    subtitle_placeholder.left = Inches(0.55)
    subtitle_placeholder.top = Inches(1.95)
    subtitle_placeholder.width = Inches(6.3)
    subtitle_placeholder.height = Inches(1.2)
    subtitle_box = subtitle_placeholder.text_frame
    subtitle_box.clear()
    subtitle_box.word_wrap = True
    subtitle_box.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    para = subtitle_box.paragraphs[0]
    para.text = "CloudCart E-commerce Platform"
    para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    run = para.runs[0]
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = BODY_BLACK

    intro_box = slide.shapes.add_textbox(Inches(0.55), Inches(2.65), Inches(6.2), Inches(1.65))
    intro_tf = intro_box.text_frame
    intro_tf.word_wrap = True
    intro_tf.text = (
        "A production-style cloud e-commerce application with secure authentication, "
        "catalog management, persistent cart, and order tracking."
    )
    intro_run = intro_tf.paragraphs[0].runs[0]
    intro_run.font.size = Pt(18)
    intro_run.font.color.rgb = BODY_BLACK

    table = slide.shapes.add_table(rows=4, cols=2, left=Inches(6.95), top=Inches(1.95), width=Inches(5.15), height=Inches(2.45)).table
    table.cell(0, 0).text = "Team Member"
    table.cell(0, 1).text = "USN / ID"
    rows = [("shreya", "258"), ("ravi", "259"), ("prajwal", "250")]
    for idx, (name, usn) in enumerate(rows, start=1):
        table.cell(idx, 0).text = name
        table.cell(idx, 1).text = usn

    for row_idx in range(4):
        for col_idx in range(2):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(232, 241, 255) if row_idx == 0 else RGBColor(250, 253, 255)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            for r in p.runs:
                r.font.size = Pt(15)
                r.font.bold = row_idx == 0
                r.font.color.rgb = TITLE_BLUE if row_idx == 0 else BODY_BLACK

    strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(4.55), Inches(11.55), Inches(2.25))
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(241, 247, 255)
    strip.line.color.rgb = RGBColor(190, 213, 245)
    strip_tf = strip.text_frame
    strip_tf.word_wrap = True
    strip_tf.text = (
        "Highlights\n"
        "• FastAPI + PostgreSQL architecture\n"
        "• Render-ready deployment pipeline\n"
        "• Admin operations + image upload\n"
        "• Real-world UX: multi-page storefront"
    )
    strip_tf.paragraphs[0].runs[0].font.size = Pt(17)
    strip_tf.paragraphs[0].runs[0].font.bold = True
    strip_tf.paragraphs[0].runs[0].font.color.rgb = TITLE_BLUE
    for paragraph in strip_tf.paragraphs[1:]:
        if paragraph.runs:
            paragraph.runs[0].font.size = Pt(15)
            paragraph.runs[0].font.color.rgb = BODY_BLACK


def add_content_slide(prs: Presentation, heading: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_theme(slide)
    add_logo(slide)

    title_shape = slide.shapes.title
    title_shape.left = Inches(0.55)
    title_shape.top = Inches(1.0)
    title_shape.width = Inches(9.0)
    title_shape.height = Inches(0.8)
    style_title(title_shape.text_frame, heading, size=30)

    body_shape = slide.placeholders[1]
    body_shape.left = Inches(0.55)
    body_shape.top = Inches(1.95)
    body_shape.width = Inches(8.3)
    body_shape.height = Inches(4.65)
    body = body_shape.text_frame
    body.clear()
    body.word_wrap = True
    for idx, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        paragraph.text = f"• {bullet}"
        paragraph.level = 0
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(18)
            run.font.bold = False
            run.font.color.rgb = BODY_BLACK

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.15), Inches(1.95), Inches(2.95), Inches(4.65))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(241, 247, 255)
    panel.line.color.rgb = RGBColor(190, 213, 245)
    panel_tf = panel.text_frame
    panel_tf.word_wrap = True
    panel_tf.text = (
        "Value Delivered\n"
        "1. Scalability\n"
        "2. Security\n"
        "3. Maintainability\n"
        "4. User Experience\n"
        "5. Deployment Readiness"
    )
    panel_tf.paragraphs[0].runs[0].font.size = Pt(16)
    panel_tf.paragraphs[0].runs[0].font.bold = True
    panel_tf.paragraphs[0].runs[0].font.color.rgb = TITLE_BLUE
    for paragraph in panel_tf.paragraphs[1:]:
        if paragraph.runs:
            paragraph.runs[0].font.size = Pt(14)
            paragraph.runs[0].font.color.rgb = BODY_BLACK


def add_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_theme(slide)
    add_logo(slide)

    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.65), Inches(11.2), Inches(2.8))
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    run = p.runs[0]
    run.font.size = Pt(64)
    run.font.bold = True
    run.font.color.rgb = TITLE_BLUE

    p2 = tf.add_paragraph()
    p2.text = "CloudCart - Production-Style E-commerce PaaS on Render"
    p2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    for r in p2.runs:
        r.font.size = Pt(30)
        r.font.bold = True
        r.font.color.rgb = TITLE_BLUE

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.45), Inches(4.6), Inches(9.9), Inches(2.0))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(241, 247, 255)
    card.line.color.rgb = RGBColor(190, 213, 245)
    card_tf = card.text_frame
    card_tf.word_wrap = True
    card_tf.text = (
        "Team\n"
        "shreya 258   |   ravi 259   |   prajwal 250\n"
        "KLE TECHNOLOGICAL UNIVERSITY"
    )
    card_tf.paragraphs[0].runs[0].font.size = Pt(18)
    card_tf.paragraphs[0].runs[0].font.bold = True
    card_tf.paragraphs[0].runs[0].font.color.rgb = TITLE_BLUE
    for paragraph in card_tf.paragraphs[1:]:
        if paragraph.runs:
            paragraph.runs[0].font.size = Pt(16)
            paragraph.runs[0].font.bold = True
            paragraph.runs[0].font.color.rgb = BODY_BLACK


def create_presentation() -> Path:
    prs = Presentation()

    add_title_slide(
        prs,
        "KLE TECHNOLOGICAL UNIVERSITY",
    )
    add_content_slide(
        prs,
        "Project Overview",
        [
            "CloudCart is a full-stack shopping application with real e-commerce workflows.",
            "Built with FastAPI, PostgreSQL, and a browser-based frontend.",
            "Designed for deployment on Render with managed cloud database support.",
        ],
    )
    add_content_slide(
        prs,
        "Architecture & Tech Stack",
        [
            "Backend: FastAPI with modular routers, schema validation, and JWT security.",
            "Database: PostgreSQL via SQLAlchemy ORM and relational table modeling.",
            "Frontend: Responsive HTML/CSS/JS dashboard consuming REST APIs.",
        ],
    )
    add_content_slide(
        prs,
        "Authentication & User Features",
        [
            "Signup/Login implemented with hashed passwords and JWT access tokens.",
            "User profile update endpoint supports contact and address management.",
            "Order history endpoint returns detailed item-level purchase records.",
        ],
    )
    add_content_slide(
        prs,
        "Product Discovery & Search",
        [
            "Public product listing endpoint with category and keyword search.",
            "Filter by price ranges and sort by price, rating, or newest items.",
            "Admin can create, edit, and delete products dynamically.",
        ],
    )
    add_content_slide(
        prs,
        "Cart & Checkout",
        [
            "Persistent cart table stores product quantities per user in database.",
            "Checkout converts cart to order with inventory deduction and totals.",
            "Order status lifecycle supports pending, shipped, and delivered.",
        ],
    )
    add_content_slide(
        prs,
        "Admin Dashboard",
        [
            "Admin-only routes secured via role-based dependency checks.",
            "Dashboard provides users, products, orders, and revenue analytics.",
            "Admin can track all orders and update shipment status quickly.",
        ],
    )
    add_content_slide(
        prs,
        "Advanced Feature: Image Upload",
        [
            "Product image upload endpoint supports jpg/jpeg/png/webp formats.",
            "Size validation and secure unique filenames prevent collisions.",
            "Uploaded images are served from static storage for product cards.",
        ],
    )
    add_content_slide(
        prs,
        "Deployment on Render",
        [
            "render.yaml defines web service, PostgreSQL, env vars, and start command.",
            "Health endpoint enables uptime checks and deployment monitoring.",
            "Environment-based settings allow smooth dev-to-production promotion.",
        ],
    )
    add_closing_slide(prs)

    pptx_path = Path("CloudCart_Presentation.pptx")
    prs.save(pptx_path)
    return pptx_path


def export_pdf(pptx_path: Path) -> Optional[Path]:
    output_dir = pptx_path.parent.resolve()
    pdf_path = output_dir / "CloudCart_Presentation.pdf"
    for command in [
        ["soffice", "--headless", "--convert-to", "pdf", str(pptx_path), "--outdir", str(output_dir)],
        ["libreoffice", "--headless", "--convert-to", "pdf", str(pptx_path), "--outdir", str(output_dir)],
    ]:
        try:
            subprocess.run(command, check=True, capture_output=True, text=True)
            if pdf_path.exists():
                return pdf_path
        except (subprocess.CalledProcessError, FileNotFoundError):
            continue
    return create_pdf_with_reportlab(output_dir / "CloudCart_Presentation.pdf")


def create_pdf_with_reportlab(pdf_path: Path) -> Optional[Path]:
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import landscape, letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
    except ImportError:
        return None

    doc = SimpleDocTemplate(
        str(pdf_path),
        pagesize=landscape(letter),
        leftMargin=0.5 * inch,
        rightMargin=0.5 * inch,
        topMargin=0.6 * inch,
        bottomMargin=0.6 * inch,
    )
    styles = getSampleStyleSheet()
    title_style = styles["Heading1"]
    title_style.textColor = colors.HexColor("#275FAA")
    body_style = styles["BodyText"]
    body_style.textColor = colors.black
    body_style.fontSize = 17
    body_style.leading = 24

    slide_data = [
        ("KLE TECHNOLOGICAL UNIVERSITY", ["CloudCart E-commerce Platform"]),
        (
            "Project Overview",
            [
                "CloudCart is a full-stack shopping application with real e-commerce workflows.",
                "Built with FastAPI, PostgreSQL, and a browser-based frontend.",
                "Designed for deployment on Render with managed cloud database support.",
            ],
        ),
        (
            "Architecture & Tech Stack",
            [
                "Backend: FastAPI with modular routers, schema validation, and JWT security.",
                "Database: PostgreSQL via SQLAlchemy ORM and relational table modeling.",
                "Frontend: Responsive HTML/CSS/JS dashboard consuming REST APIs.",
            ],
        ),
        (
            "Authentication & User Features",
            [
                "Signup/Login implemented with hashed passwords and JWT access tokens.",
                "User profile update endpoint supports contact and address management.",
                "Order history endpoint returns detailed item-level purchase records.",
            ],
        ),
        (
            "Product Discovery & Search",
            [
                "Public product listing endpoint with category and keyword search.",
                "Filter by price ranges and sort by price, rating, or newest items.",
                "Admin can create, edit, and delete products dynamically.",
            ],
        ),
        (
            "Cart & Checkout",
            [
                "Persistent cart table stores product quantities per user in database.",
                "Checkout converts cart to order with inventory deduction and totals.",
                "Order status lifecycle supports pending, shipped, and delivered.",
            ],
        ),
        (
            "Admin Dashboard",
            [
                "Admin-only routes secured via role-based dependency checks.",
                "Dashboard provides users, products, orders, and revenue analytics.",
                "Admin can track all orders and update shipment status quickly.",
            ],
        ),
        (
            "Advanced Feature: Image Upload",
            [
                "Product image upload endpoint supports jpg/jpeg/png/webp formats.",
                "Size validation and secure unique filenames prevent collisions.",
                "Uploaded images are served from static storage for product cards.",
            ],
        ),
        (
            "Deployment on Render",
            [
                "render.yaml defines web service, PostgreSQL, env vars, and start command.",
                "Health endpoint enables uptime checks and deployment monitoring.",
                "Environment-based settings allow smooth dev-to-production promotion.",
            ],
        ),
        ("Thank You", ["CloudCart - Production-Style E-commerce PaaS on Render"]),
    ]

    story = []
    for idx, (title, bullets) in enumerate(slide_data):
        story.append(Paragraph(f"<u>{title}</u>", title_style))
        story.append(Spacer(1, 0.2 * inch))
        if idx == 0:
            story.append(Paragraph(bullets[0], body_style))
            story.append(Spacer(1, 0.15 * inch))
            team_table = Table(
                [["Team Member", "USN / ID"], ["shreya", "258"], ["ravi", "259"], ["prajwal", "250"]],
                colWidths=[2.4 * inch, 1.6 * inch],
            )
            team_table.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E8F1FF")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#275FAA")),
                        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                        ("GRID", (0, 0), (-1, -1), 1, colors.HexColor("#B7CBEF")),
                        ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#FAFDFF")),
                    ]
                )
            )
            story.append(team_table)
        else:
            for bullet in bullets:
                story.append(Paragraph(f"• {bullet}", body_style))
                story.append(Spacer(1, 0.08 * inch))
        story.append(Spacer(1, 0.35 * inch))
        if idx < len(slide_data) - 1:
            story.append(PageBreak())

    def draw_page(canvas, _doc):
        canvas.saveState()
        canvas.setStrokeColorRGB(0.77, 0.12, 0.22)
        canvas.setLineWidth(1)
        canvas.line(30, 535, 770, 535)
        if LOGO_PATH.exists():
            canvas.drawImage(str(LOGO_PATH), 560, 545, width=210, height=42, preserveAspectRatio=True)
        canvas.restoreState()

    doc.build(story, onFirstPage=draw_page, onLaterPages=draw_page)
    return pdf_path if pdf_path.exists() else None


if __name__ == "__main__":
    pptx_output = create_presentation()
    pdf_output = export_pdf(pptx_output)
    print(f"Created {pptx_output.name}")
    if pdf_output:
        print(f"Created {pdf_output.name}")
    else:
        print("PDF export tool not found. Install LibreOffice to enable automatic PDF export.")
